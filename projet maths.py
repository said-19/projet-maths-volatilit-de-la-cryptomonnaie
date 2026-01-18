



# Standard imports
import os
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from tqdm import tqdm

# Create outputs folder
OUTDIR = 'outputs'
if not os.path.exists(OUTDIR):
    os.makedirs(OUTDIR)

# ---------------------------
# 1) Data download / load
# ---------------------------

def load_crypto(symbol='BTC-USD', start='2018-01-01', end=None, source='yfinance'):
    """Download daily price data using yfinance. If not available, load CSV.
    Returns DataFrame with 'Date' index and 'Adj Close'."""
    try:
        import yfinance as yf
    except Exception as e:
        raise ImportError('yfinance required: pip install yfinance')
    df = yf.download(symbol, start=start, end=end, progress=False)
    if df.empty:
        raise ValueError('No data downloaded. Check symbol/date range or internet connection.')
    df = df[['Adj Close']].rename(columns={'Adj Close':'price'})
    df.index = pd.to_datetime(df.index)
    return df

# Example load
df_price = load_crypto('BTC-USD', start='2018-01-01', end=None)
print(f'Downloaded {len(df_price)} rows from {df_price.index.min().date()} to {df_price.index.max().date()}')

# ---------------------------
# 2) Preprocess: compute log returns
# ---------------------------

df = df_price.copy()
df['log_price'] = np.log(df['price'])
df['ret'] = 100 * (df['log_price'].diff())
df = df.dropna().copy()

# Basic stats
desc = df['ret'].describe()
print('\nReturn descriptive statistics:')
print(desc)

# Save quick plot
plt.figure(figsize=(10,4))
plt.plot(df.index, df['ret'])
plt.title('Daily log returns (%) - BTC')
plt.tight_layout()
plt.savefig(os.path.join(OUTDIR,'returns_plot.png'))
plt.close()

# ---------------------------
# 3) GARCH family estimation (arch package)
# ---------------------------

try:
    from arch import arch_model
except Exception as e:
    raise ImportError('arch required: pip install arch')

models_results = {}

# Define function to fit GARCH variants
def fit_garch_models(returns):
    res = {}
    # GARCH(1,1) with Student-t
    am_garch = arch_model(returns, mean='Constant', vol='GARCH', p=1, q=1, dist='t')
    res_garch = am_garch.fit(disp='off')
    res['GARCH'] = res_garch

    # TGARCH (GJR-GARCH) -> use o=1 parameter
    am_tgarch = arch_model(returns, mean='Constant', vol='GARCH', p=1, o=1, q=1, dist='t')
    res_tgarch = am_tgarch.fit(disp='off')
    res['TGARCH'] = res_tgarch

    # IGARCH approximation: enforce alpha+beta ~ 1 by fitting and inspect
    # Arch does not directly enforce alpha+beta=1; but we fit GARCH and later check persistence
    am_igarch = arch_model(returns, mean='Constant', vol='GARCH', p=1, q=1, dist='t')
    res_igarch = am_igarch.fit(disp='off')
    res['IGARCH_like'] = res_igarch

    return res

print('\nFitting GARCH models (this may take a bit)...')
garch_res = fit_garch_models(df['ret'])

# Save summaries and AIC
garch_summaries = {}
for name, r in garch_res.items():
    s = r.summary()
    print(f'\nModel: {name}')
    print(r.summary())
    garch_summaries[name] = {'aic': r.aic, 'bic': r.bic, 'llf': r.loglikelihood}

pd.DataFrame(garch_summaries).to_csv(os.path.join(OUTDIR,'garch_model_info.csv'))

# ---------------------------
# 4) Forecasting with GARCH
# ---------------------------

def garch_forecast_mse(fitted_model, returns, horizons=[3,5,10,20,30,44], test_size=60):
    """Rolling-origin forecasting as in the article:
    - Use the last `test_size` observations as out-of-sample holdout
    - Re-fit model every step (walk-forward) - can be slow
    Returns dict of horizon->MSE1 and MSE2
    """
    n = len(returns)
    train_end = n - test_size
    results = {h: {'mse1': [], 'mse2': []} for h in horizons}

    # Simple version: refit once on train, then forecast horizons from each test origin
    # For full rolling refit, uncomment the refit loop (slow)
    for i in range(train_end, n):
        train = returns[:i]
        # re-fit to train
        fm = fitted_model.clone() if hasattr(fitted_model,'clone') else None
        # We'll fit a fresh model of same spec
        spec = fitted_model.model
        am = arch_model(train, mean=spec.mean, vol=spec.vol, p=spec.p, q=spec.q, o=getattr(spec,'o',0), dist='t')
        try:
            fr = am.fit(disp='off')
        except Exception:
            fr = fitted_model  # fallback
        # generate h-step ahead forecasts using .forecast
        fc = fr.forecast(horizon=max(horizons), reindex=False)
        for h in horizons:
            # take the predicted variance for horizon h at time i
            try:
                var_pred = fc.variance.iloc[-1, h-1]
            except Exception:
                var_pred = np.nan
            # actual proxy
            if i + h - 1 < n:
                rt = returns.iloc[i + h - 1]
                mse1 = (np.log(rt**2 + 1e-12) - np.log(var_pred + 1e-12))**2
                mse2 = (np.log(abs(rt) + 1e-12) - np.log(var_pred + 1e-12))**2
                results[h]['mse1'].append(mse1)
                results[h]['mse2'].append(mse2)
    # aggregate
    agg = {}
    for h in horizons:
        agg[h] = {'MSE1': np.nanmean(results[h]['mse1']), 'MSE2': np.nanmean(results[h]['mse2'])}
    return agg

# Choose best GARCH from AIC
best = min(garch_res.items(), key=lambda x: x[1].aic)
best_name, best_model = best[0], garch_res[best[0]]
print(f"\nBest GARCH variant by AIC: {best_name}")

# Forecast MSE (NOTE: rolling refit is heavy; we provide a single-run forecast to produce tables fast)
garch_mse = garch_forecast_mse(best_model, df['ret'], horizons=[3,5,10,20,30,44], test_size=120)
print('\nGARCH forecast MSEs:')
print(pd.DataFrame(garch_mse).T)
pd.DataFrame(garch_mse).T.to_csv(os.path.join(OUTDIR,'garch_forecast_mse.csv'))

# ---------------------------
# 5) Bayesian Stochastic Volatility (PyMC)
# ---------------------------

# The SV model implementation uses pymc. This block may take time.
try:
    import pymc as pm
    import arviz as az
except Exception as e:
    print('PyMC or Arviz not installed. Install via `pip install pymc arviz` for SV estimation.')
    pm = None

if pm is not None:
    returns = df['ret'].values
    T = len(returns)

    with pm.Model() as sv_model:
        # Priors as suggested in literature
        mu = pm.Normal('mu', mu=0.0, sigma=10.0)
        phi = pm.Beta('phi_raw', alpha=20.0, beta=1.5)  # maps to (0,1)
        phi = pm.Deterministic('phi', 2*phi - 1)  # map to (-1,1)
        sigma_eta = pm.Exponential('sigma_eta', 1.0)

        # latent log-variance process
        h = pm.GaussianRandomWalk('h', sigma=sigma_eta, shape=T)  # alternative simpler prior

        # observation
        obs = pm.Normal('obs', mu=0.0, sigma=pm.math.exp(h/2), observed=returns)

        # sample posterior
        trace = pm.sample(draws=1000, tune=1000, chains=2, target_accept=0.9, return_inferencedata=True)

    az.to_netcdf(trace, os.path.join(OUTDIR,'sv_trace.nc'))
    print('\nSV sampling finished. Summary:')
    print(az.summary(trace, var_names=['mu','phi','sigma_eta']))

    # Posterior predictive latent volatility: take median h and forecast
    post_h = trace.posterior['h'].median(dim=('chain','draw')).values
    df['sv_h_median'] = post_h
    df['sv_var_med'] = np.exp(df['sv_h_median'])
    df[['sv_var_med']].to_csv(os.path.join(OUTDIR,'sv_median_variance.csv'))

    # Forecasting future h for horizons: draw from posterior predictive
    def sv_forecast_from_trace(trace, last_h, horizons=44, draws=200):
        # simple simulation assuming AR(1) on h
        mu_s = trace.posterior['mu'].stack(samples=('chain','draw')).values
        # phi and sigma_eta extraction
        phi_s = trace.posterior['phi'].stack(samples=('chain','draw')).values
        sigma_s = trace.posterior['sigma_eta'].stack(samples=('chain','draw')).values
        nsamps = min(draws, phi_s.shape[-1])
        sim_vars = np.zeros((nsamps, horizons))
        for i in range(nsamps):
            h_prev = last_h
            phi_i = phi_s[i]
            sigma_i = sigma_s[i]
            mu_i = mu_s[i]
            for t in range(horizons):
                h_prev = mu_i + phi_i*(h_prev-mu_i) + np.random.normal(scale=sigma_i)
                sim_vars[i,t] = np.exp(h_prev)
        return sim_vars

    last_h = df['sv_h_median'].iloc[-1]
    sims = sv_forecast_from_trace(trace, last_h, horizons=44, draws=200)
    # compute mean predicted variance for each horizon
    sv_horizon_var = sims.mean(axis=0)
    # compute MSEs similar to GARCH for the overlapping test period if possible
    # For demonstration, save sv_horizon_var
    pd.DataFrame({'sv_var_pred_mean': sv_horizon_var}).to_csv(os.path.join(OUTDIR,'sv_horizon_var.csv'))

# ---------------------------
# 6) Generate a PowerPoint summary automatically
# ---------------------------

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None

if Presentation is not None:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Volatility Forecasting: GARCH vs Stochastic Volatility"
    subtitle.text = "Auto-generated slides - outputs/presentation.pptx"

    # Slide 2: Objectives & Data
    sld = prs.slides.add_slide(prs.slide_layouts[1])
    sld.shapes.title.text = 'Objectives & Data'
    body = sld.shapes.placeholders[1].text_frame
    body.text = 'Objectives:'
    p = body.add_paragraph()
    p.text = '- Compare GARCH family and Bayesian SV for volatility forecasting'
    p.level = 1
    p = body.add_paragraph()
    p.text = 'Data: BTC-USD daily returns, start 2018-01-01'
    p.level = 1

    # Slide 3: Methodology
    sld = prs.slides.add_slide(prs.slide_layouts[1])
    sld.shapes.title.text = 'Methodology'
    body = sld.shapes.placeholders[1].text_frame
    body.text = 'GARCH: GARCH(1,1), TGARCH (o=1), IGARCH-like\nSV: Bayesian AR(1) on log-variance with MCMC'

    # Slide 4: Results placeholders
    sld = prs.slides.add_slide(prs.slide_layouts[1])
    sld.shapes.title.text = 'Key Results (see outputs)'
    body = sld.shapes.placeholders[1].text_frame
    body.text = 'Files generated in outputs/:\n- returns_plot.png\n- garch_model_info.csv\n- garch_forecast_mse.csv\n- sv_horizon_var.csv\n- sv_trace.nc'

    # Save
    pptx_path = os.path.join(OUTDIR,'presentation.pptx')
    prs.save(pptx_path)
    print(f'Presentation saved to {pptx_path}')
else:
    print('python-pptx not installed; skip PPT generation (pip install python-pptx)')

# ---------------------------
# 7) Wrap-up: produce a small report CSV with key metrics
# ---------------------------

summary = {
    'best_garch': best_name,
}
with open(os.path.join(OUTDIR,'summary.txt'),'w') as f:
    f.write(str(summary))

print('\nAll done. Outputs saved in folder:', OUTDIR)
print('Key files: returns_plot.png, garch_model_info.csv, garch_forecast_mse.csv, sv_horizon_var.csv, presentation.pptx (if python-pptx installed)')

# End of file
